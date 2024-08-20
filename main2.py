from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Artificial Intelligence: An Overview"
subtitle.text = "Understanding the Basics, Applications, and Future of AI\nPresenter Name\nDate of Presentation"

# Slide 2: Introduction to AI
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "What is Artificial Intelligence?"
content.text = ("Definition: AI is the simulation of human intelligence in machines that are programmed to think and learn like humans.\n"
                "Key Concepts: Automation, Machine Learning, Neural Networks")

# Slide 3: History of AI
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]
title.text = "The Evolution of AI"
content.text = ("Early Beginnings: 1950s - Alan Turing's 'Computing Machinery and Intelligence'\n"
                "First AI Programs: 1956 - Dartmouth Conference\n"
                "AI Winters: Periods of reduced funding and interest\n"
                "Modern AI: 21st Century - Advances in computational power and big data")

# Slide 4: Types of AI
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]
title.text = "Categories of AI"
content.text = ("Narrow AI (Weak AI): AI systems designed for specific tasks (e.g., voice assistants, recommendation systems)\n"
                "General AI (Strong AI): Hypothetical AI with human-like cognitive abilities\n"
                "Super AI: AI surpassing human intelligence (future concept)")

# Slide 5: Machine Learning and Deep Learning
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]
title.text = "Core Technologies in AI"
content.text = ("Machine Learning: Algorithms that allow machines to learn from data (e.g., supervised, unsupervised learning)\n"
                "Deep Learning: Subset of ML using neural networks with many layers (e.g., image and speech recognition)")

# Slide 6: Applications of AI
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]
title.text = "Real-World Applications"
content.text = ("Healthcare: Disease diagnosis, personalized treatment plans\n"
                "Finance: Fraud detection, algorithmic trading\n"
                "Transportation: Autonomous vehicles, traffic management\n"
                "Retail: Personalized shopping experiences, inventory management")

# Slide 7: Ethical Considerations in AI
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.placeholders[1]
title.text = "Ethical and Social Implications"
content.text = ("Bias and Fairness: Ensuring AI systems are unbiased\n"
                "Privacy: Handling of personal data\n"
                "Job Displacement: Impact on employment\n"
                "Accountability: Responsibility for AI decisions")

# Slide 8: Future of AI
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.placeholders[1]
title.text = "What Lies Ahead?"
content.text = ("Technological Advances: Improvements in algorithms, hardware\n"
                "New Applications: AI in new industries (e.g., education, law)\n"
                "Regulations and Policies: Developing frameworks to govern AI use")

# Slide 9: AI Tools and Technologies
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.placeholders[1]
title.text = "Key AI Tools"
content.text = ("Programming Languages: Python, R\n"
                "Frameworks: TensorFlow, PyTorch, Keras\n"
                "Platforms: IBM Watson, Google AI, Amazon AI")

# Slide 10: Conclusion
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.placeholders[1]
title.text = "Summary and Closing"
content.text = ("Recap: Brief overview of key points\n"
                "Future Prospects: AI’s potential to transform industries\n"
                "Q&A: Open the floor for questions")

# Slide 11: References
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_11.shapes.title
content = slide_11.placeholders[1]
title.text = "Sources and Further Reading"
content.text = "List of books, articles, and websites referenced in the presentation"

# Slide 12: Thank You
slide_12 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_12.shapes.title
content = slide_12.placeholders[1]
title.text = "Thank You for Your Attention!"
content.text = "Contact Information: Your email or other contact details"

# Save the presentation
file_path = "/mnt/data/AI_Presentation.pptx"
prs.save(file_path)

file_path
